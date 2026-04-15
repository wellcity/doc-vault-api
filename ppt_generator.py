import sys
sys.path.append(".")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import io
import base64
from config import OUTPUT_DIR


def generate_ppt(chunks: list[dict], output_name: str = "report", include_images: bool = True) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 空白版面

    for chunk in chunks:
        slide = prs.slides.add_slide(blank_layout)

        # === 標題列 ===
        title_left = Inches(0.3)
        title_top = Inches(0.2)
        title_width = Inches(12.73)
        title_height = Inches(0.6)

        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = title_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"[{chunk.get('file_type', '').upper()}] {chunk.get('source_file', '')} — 第 {chunk.get('page', 0)} 頁"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

        # === 文字內容區 ===
        text_left = Inches(0.3)
        text_top = Inches(1.0)
        text_width = Inches(7.0)
        text_height = Inches(6.0)

        text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = chunk.get("text", "")
        p.font.size = Pt(12)

        # === 圖片區 ===
        if include_images:
            image_paths = chunk.get("image_paths", [])
            if image_paths:
                img_start_x = Inches(7.5)
                img_start_y = Inches(1.0)
                max_img_width = Inches(5.5)
                max_img_height = Inches(6.0)

                for idx, img_path in enumerate(image_paths[:4]):  # 最多4張
                    try:
                        path = Path(img_path)
                        if not path.exists():
                            continue

                        # Grid 排列
                        row = idx // 2
                        col = idx % 2
                        x = img_start_x + Inches(col * 2.7)
                        y = img_start_y + Inches(row * 3.0)

                        slide.shapes.add_picture(str(path), x, y, width=max_img_width / 2, height=max_img_height / 2)
                    except Exception as e:
                        print(f"Warning: cannot add image {img_path}: {e}")

    # 儲存到 BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


# =============================================================================
# 大綱生成模式
# =============================================================================

def generate_ppt_from_outline(outline: dict, output_name: str = "presentation") -> bytes:
    """
    根據大綱結構生成 PPT。

    Args:
        outline: 大綱結構，格式如下：
            {
                "title": "報告標題",
                "subtitle": "副標題（選填）",
                "author": "作者（選填）",
                "slides": [
                    {
                        "title": "第一頁標題",
                        "bullets": ["項目一", "項目二", "項目三"],  # 選填
                        "notes": "備註（選填）"
                    },
                    ...
                ]
            }
        output_name: 輸出檔名（不含副檔名）

    Returns:
        PPT 檔案內容（bytes）
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 全域配色
    COLOR_TITLE = RGBColor(0x1F, 0x49, 0x7D)   # 深藍
    COLOR_BULLET = RGBColor(0x26, 0x26, 0x26)  # 深灰
    COLOR_ACCENT = RGBColor(0x2E, 0x86, 0xAB)  # 亮藍

    # --- 封面頁 ---
    cover_layout = prs.slide_layouts[6]  # 空白版面
    cover = prs.slides.add_slide(cover_layout)

    # 全高深藍色矩形（頂部區塊）
    rect_top = cover.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(13.33), Inches(3.5)
    )
    rect_top.fill.solid()
    rect_top.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    rect_top.line.fill.background()

    # 報告標題
    title_box = cover.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.73), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = outline.get("title", "簡報")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 副標題
    subtitle = outline.get("subtitle")
    if subtitle:
        sub_box = cover.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.73), Inches(0.6))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)

    # 作者 + 日期（底部）
    author = outline.get("author")
    import datetime
    date_str = datetime.date.today().strftime("%Y/%m/%d")

    footer_box = cover.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.73), Inches(0.5))
    tf3 = footer_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"{author}  |  {date_str}" if author else date_str
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # --- 內容頁 ---
    content_layout = prs.slide_layouts[6]
    slides_data = outline.get("slides", [])

    for slide_data in slides_data:
        slide = prs.slides.add_slide(content_layout)

        slide_title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        notes = slide_data.get("notes")

        # 左側標題豎條
        rect_bar = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0),
            Inches(0.15), Inches(7.5)
        )
        rect_bar.fill.solid()
        rect_bar.fill.fore_color.rgb = COLOR_ACCENT
        rect_bar.line.fill.background()

        # 頁面標題
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_TITLE

        # 分隔線
        line = slide.shapes.add_shape(
            1,
            Inches(0.4), Inches(1.15),
            Inches(12.5), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.fill.background()

        # 重点项目
        if bullets:
            bullet_top = Inches(1.4)
            for idx, bullet in enumerate(bullets):
                bullet_box = slide.shapes.add_textbox(
                    Inches(0.6), bullet_top,
                    Inches(12.0), Inches(0.65)
                )
                tf = bullet_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]

                # 圆点编号
                p.text = f"{idx + 1}.  {bullet}"
                p.font.size = Pt(18)
                p.font.color.rgb = COLOR_BULLET
                p.space_after = Pt(6)

                bullet_top += Inches(0.65)

        # 備註
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
