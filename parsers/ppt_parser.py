import sys
sys.path.append(".")

import uuid
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
from config import IMAGE_DIR
from datetime import datetime, timezone
import json


def parse_pptx(file_path: str, file_id: str, metadata: dict | None = None) -> tuple[list[dict], list[str]]:
    """
    解析 PowerPoint (.pptx) 檔案，回傳 (chunks, image_paths)
    每張投影片為一個 chunk
    """
    chunks = []
    image_paths = []
    metadata = metadata or {}

    prs = Presentation(file_path)
    base_name = Path(file_path).stem

    img_dir = Path(IMAGE_DIR) / file_id
    img_dir.mkdir(parents=True, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = ""

        # 標題
        if slide.shapes.title:
            slide_text += f"# {slide.shapes.title.text}\n"

        # 內文
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text += text + "\n"

        if not slide_text.strip():
            continue

        # 萃取圖片
        slide_image_paths = []
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    img = shape.image
                    img_bytes = img.blob
                    img_ext = img.ext
                    img_name = f"{base_name}_slide{slide_num}_{len(slide_image_paths)+1}.{img_ext}"
                    img_path = str(img_dir / img_name)
                    with open(img_path, "wb") as f:
                        f.write(img_bytes)
                    image_paths.append(img_path)
                    slide_image_paths.append(img_path)
                except Exception:
                    pass

        # 每張投影片一個 chunk
        chunk = {
            "chunk_id": str(uuid.uuid4()),
            "file_id": file_id,
            "source_file": Path(file_path).name,
            "file_type": "ppt",
            "page": slide_num,
            "chunk_index": 0,
            "text": slide_text.strip(),
            "image_paths": slide_image_paths,
            "metadata_json": json.dumps(metadata),
            "embedding": None,
            "created_at": datetime.now(timezone.utc).isoformat(),
        }
        chunks.append(chunk)

    return chunks, image_paths
